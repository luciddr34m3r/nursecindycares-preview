"""
Rebuild: Beyond Burnout — Allied Health Symposium
Nurse Cindy (Cindy Barnard BSN, RN, OCN)

Full content from original deck, completely redesigned with Nurse Cindy brand.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os, copy

# ── Brand ─────────────────────────────────────────────────────────────────────
PLUM     = RGBColor(0x2C, 0x2A, 0x4A)
PLUM2    = RGBColor(0x3D, 0x3A, 0x6E)   # lighter plum for panels
CORAL    = RGBColor(0xE8, 0x89, 0x6A)
CREAM    = RGBColor(0xFA, 0xF7, 0xF4)
BLUSH    = RGBColor(0xF2, 0xD4, 0xCC)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
CHARCOAL = RGBColor(0x3D, 0x3D, 0x3D)
MUTED    = RGBColor(0x8A, 0x88, 0xA8)
PALE     = RGBColor(0xC8, 0xC6, 0xD8)

IMG  = "/Users/CurtBarnard/code/nursecindycares.com/assets/images/deck/"
CINDY_HERO = "/Users/CurtBarnard/code/nursecindycares.com/assets/images/cindy-hero.png"
CINDY_ABOUT = "/Users/CurtBarnard/code/nursecindycares.com/assets/images/cindy-about.jpg"

OUT  = "/Users/CurtBarnard/code/nursecindycares.com/assets/NurseCindy-BeyondBurnout.pptx"

W  = Inches(13.33)
H  = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ── Helpers ───────────────────────────────────────────────────────────────────

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, color):
    shp = slide.shapes.add_shape(1, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp

def pic(slide, path, l, t, w, h):
    if os.path.exists(path):
        return slide.shapes.add_picture(path, l, t, w, h)
    return None

def tb(slide, l, t, w, h, text, size, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False, wrap=True, spacing=None):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    if spacing:
        p.space_after = Pt(spacing)
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name  = "Calibri"
    return box

def multiline(slide, l, t, w, h, lines, size, bold=False, color=WHITE,
              align=PP_ALIGN.LEFT, italic=False, line_spacing=None):
    """Multiple paragraphs in one textbox."""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    for idx, (text, sz, bld, col) in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.space_after = Pt(line_spacing)
        run = p.add_run()
        run.text = text
        run.font.size  = Pt(sz)
        run.font.bold  = bld
        run.font.italic = italic
        run.font.color.rgb = col
        run.font.name  = "Calibri"
    return box

def eyebrow(slide, text, top=Inches(0.55), color=CORAL):
    tb(slide, Inches(1), top, Inches(11.33), Inches(0.4),
       text.upper(), 10, bold=True, color=color, align=PP_ALIGN.CENTER)

def side_bar(slide, color=CORAL, width=Inches(0.22)):
    rect(slide, 0, 0, width, H, color)

def top_bar(slide, color=CORAL, height=Inches(0.15)):
    rect(slide, 0, 0, W, height, color)

def bullet_rows(slide, items, x, y, w, size=11.5, icon_color=CORAL, text_color=WHITE,
                icon="•", row_h=Inches(0.55)):
    for i, item in enumerate(items):
        tb(slide, x, y + i * row_h, Inches(0.3), row_h, icon, size + 2, bold=True, color=icon_color)
        tb(slide, x + Inches(0.3), y + i * row_h, w - Inches(0.3), row_h, item, size, color=text_color)

def divider_line(slide, l, t, w, color=CORAL, h=Pt(2.5)):
    rect(slide, l, t, w, h, color)

def stat_box(slide, x, y, w, h, number, label, bg_color=PLUM2, num_color=CORAL):
    rect(slide, x, y, w, h, bg_color)
    tb(slide, x, y + Inches(0.15), w, Inches(0.9),
       number, 40, bold=True, color=num_color, align=PP_ALIGN.CENTER)
    tb(slide, x, y + Inches(1.0), w, Inches(0.6),
       label, 12, color=WHITE, align=PP_ALIGN.CENTER)

def numbered_row(slide, x, y, number, title, desc,
                 num_color=CORAL, title_color=WHITE, desc_color=PALE,
                 w=Inches(12), title_size=14, desc_size=11):
    tb(slide, x, y, Inches(0.5), Inches(0.5), str(number), title_size + 2, bold=True, color=num_color)
    tb(slide, x + Inches(0.55), y, w, Inches(0.4), title, title_size, bold=True, color=title_color)
    tb(slide, x + Inches(0.55), y + Inches(0.38), w, Inches(0.55), desc, desc_size, color=desc_color)

# ── SLIDE COUNTER ─────────────────────────────────────────────────────────────
slide_n = [0]

def new_slide():
    slide_n[0] += 1
    return prs.slides.add_slide(BLANK)

def slide_num_label(slide, n=None):
    num = n or slide_n[0]
    tb(slide, W - Inches(1.2), H - Inches(0.45), Inches(1.0), Inches(0.4),
       str(num), 10, color=MUTED, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, PLUM)

side_bar(s)

# Big decorative circle top-right
circ = s.shapes.add_shape(9, Inches(10.8), Inches(-1.5), Inches(5.5), Inches(5.5))
circ.fill.solid(); circ.fill.fore_color.rgb = PLUM2; circ.line.fill.background()

# Cindy photo right side
pic(s, CINDY_HERO, Inches(9.0), Inches(0), Inches(4.33), H)

# Dark overlay on right for readability — use xml to set transparency
overlay = rect(s, Inches(9.0), 0, Inches(4.33), H, PLUM)
sp = overlay._element
spPr = sp.find(qn('p:spPr'))
solidFill = spPr.find('.//' + qn('a:solidFill'))
if solidFill is not None:
    srgb = solidFill.find(qn('a:srgbClr'))
    if srgb is None:
        srgb = etree.SubElement(solidFill, qn('a:srgbClr'))
    srgb.set('val', '2C2A4A')
    alpha_el = etree.SubElement(srgb, qn('a:alpha'))
    alpha_el.set('val', '75000')  # 75% opacity

# Eyebrow
tb(s, Inches(0.7), Inches(1.4), Inches(8), Inches(0.45),
   "ALLIED HEALTH SYMPOSIUM  ·  AUGUST 29, 2026", 10, bold=True, color=CORAL)

# Title
tb(s, Inches(0.7), Inches(2.1), Inches(8.3), Inches(1.6),
   "Beyond Burnout:", 36, bold=True, color=WHITE)
tb(s, Inches(0.7), Inches(3.5), Inches(8.3), Inches(0.9),
   "Understanding the Real\nCauses of Exhaustion", 28, bold=False, color=BLUSH, italic=True)

# Divider
divider_line(s, Inches(0.7), Inches(4.65), Inches(4))

# Speaker block
tb(s, Inches(0.7), Inches(4.85), Inches(8), Inches(0.45),
   "Cindy Barnard, BSN, RN, OCN", 16, bold=True, color=WHITE)
tb(s, Inches(0.7), Inches(5.35), Inches(8), Inches(0.35),
   "Certified 7 Types of Rest® Facilitator", 12, color=CORAL, italic=True)
tb(s, Inches(0.7), Inches(5.75), Inches(8), Inches(0.35),
   "CIONS Precision Medicine Nursing", 12, color=PALE)
tb(s, Inches(0.7), Inches(6.4), Inches(8), Inches(0.4),
   "nursecindycares.com", 11, color=MUTED)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA / INTRO
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, CREAM)
top_bar(s, PLUM)

# Right photo
pic(s, IMG + "s02-intro.png", Inches(7.8), Inches(0.15), Inches(5.33), H - Inches(0.15))
# gradient overlay on left edge of photo
grad = rect(s, Inches(7.8), Inches(0.15), Inches(1.5), H - Inches(0.15), CREAM)

eyebrow(s, "Today's Session", top=Inches(0.45), color=CORAL)

tb(s, Inches(0.8), Inches(1.1), Inches(7), Inches(1.0),
   "Discover Why Exhaustion\nIsn't Just About Sleep", 28, bold=True, color=PLUM)

tb(s, Inches(0.8), Inches(2.4), Inches(6.5), Inches(0.8),
   "A practical path to sustainable renewal that addresses the root causes of personal and professional depletion.",
   13, color=CHARCOAL)

divider_line(s, Inches(0.8), Inches(3.35), Inches(3.5), CORAL)

agenda = [
    "The burnout myth — and what WHO actually says",
    "Why time off doesn't always restore you",
    "The 7 Types of Rest® framework",
    "Common rest mistakes high performers make",
    "Micro-rest strategies and your 24-hour reset",
]
for i, item in enumerate(agenda):
    y = Inches(3.55 + i * 0.62)
    rect(s, Inches(0.8), y + Inches(0.08), Inches(0.26), Inches(0.26), CORAL)
    tb(s, Inches(1.2), y, Inches(5.8), Inches(0.55), item, 12, color=CHARCOAL)

slide_num_label(s)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE MYTH
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, PLUM)
side_bar(s)

# Left: photo half
pic(s, IMG + "s03-sleep.png", Inches(0.22), 0, Inches(5.5), H)
# Gradient over photo
grad = rect(s, Inches(4.0), 0, Inches(1.72), H, PLUM)

eyebrow(s, "The Myth", top=Inches(0.5))

tb(s, Inches(6.1), Inches(1.0), Inches(6.8), Inches(1.2),
   '"If I Could Just\nSleep More…"', 32, bold=True, color=WHITE, italic=True)

divider_line(s, Inches(6.1), Inches(2.4), Inches(4), CORAL)

tb(s, Inches(6.1), Inches(2.6), Inches(6.8), Inches(0.45),
   "Sleep Matters — But It's Not Everything", 16, bold=True, color=BLUSH)

tb(s, Inches(6.1), Inches(3.1), Inches(6.8), Inches(1.1),
   "Quality sleep is absolutely essential. But if sleep is the only form of rest you're getting, you may still be depleted in other critical dimensions.",
   12.5, color=PALE)

tb(s, Inches(6.1), Inches(4.4), Inches(6.8), Inches(1.1),
   "Exhaustion can be mental, emotional, social, sensory, creative, or spiritual. That's why \"more sleep\" helps a little — but doesn't fully restore your energy.",
   12.5, color=PALE)

divider_line(s, Inches(6.1), Inches(5.8), Inches(6.5), CORAL)
tb(s, Inches(6.1), Inches(6.0), Inches(6.8), Inches(0.5),
   "The solution isn't always more sleep. It's the right kind of rest.", 12, bold=True, color=CORAL, italic=True)

slide_num_label(s)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — WHO DEFINITION
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, WHITE)
top_bar(s, PLUM)

eyebrow(s, "Burnout Defined", color=PLUM)

tb(s, Inches(1), Inches(1.1), Inches(11.33), Inches(0.8),
   "What the WHO Actually Says", 30, bold=True, color=PLUM, align=PP_ALIGN.CENTER)

# Big quote block
rect(s, Inches(1.5), Inches(2.1), Inches(10.33), Inches(1.8), PLUM)
tb(s, Inches(1.8), Inches(2.2), Inches(9.73), Inches(1.6),
   '"Burnout is an occupational phenomenon resulting from chronic workplace stress that has not been successfully managed."',
   16, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

tb(s, Inches(1.5), Inches(4.1), Inches(10.33), Inches(0.4),
   "— World Health Organization (ICD-11)", 12, color=MUTED, align=PP_ALIGN.CENTER)

key_points = [
    ("🔵", "Specifically tied to workplace stress — not personal weakness"),
    ("🔵", "Result of unsuccessful stress management — not inevitable"),
    ("🔵", "Not classified as a medical condition — but deeply impacts health"),
    ("🔵", "Recognized globally — valid, serious, and preventable"),
]
for i, (icon, point) in enumerate(key_points):
    y = Inches(4.7 + i * 0.55)
    tb(s, Inches(1.5), y, Inches(0.4), Inches(0.5), icon, 14)
    tb(s, Inches(2.0), y, Inches(10), Inches(0.5), point, 13, color=CHARCOAL)

slide_num_label(s)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — BURNOUT VS BUSY WEEK
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, PLUM)
side_bar(s)

eyebrow(s, "Understanding the Difference")
tb(s, Inches(1), Inches(1.0), Inches(11.33), Inches(0.8),
   "Burnout vs. \"A Busy Week\"", 30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Three stage cards
stages = [
    ("STRESS SPIKE", "Temporary increase in demands with a clear endpoint", "Normal & manageable"),
    ("CHRONIC STRAIN", "Sustained high demand with no recovery window", "Warning zone"),
    ("CAPACITY SHRINKS", "Internal resources depleted over time", "Burnout territory"),
]
colors = [RGBColor(0x3D,0x3A,0x6E), RGBColor(0x4D,0x3A,0x7A), CORAL]
for i, (title, desc, sub) in enumerate(stages):
    x = Inches(0.8 + i * 4.1)
    rect(s, x, Inches(2.1), Inches(3.8), Inches(3.5), colors[i])
    tb(s, x, Inches(2.25), Inches(3.8), Inches(0.5),
       title, 11, bold=True, color=BLUSH if i < 2 else WHITE, align=PP_ALIGN.CENTER)
    divider_line(s, x + Inches(0.4), Inches(2.85), Inches(3.0),
                 CORAL if i < 2 else WHITE)
    tb(s, x + Inches(0.2), Inches(3.05), Inches(3.4), Inches(1.2),
       desc, 13, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x + Inches(0.6), Inches(4.5), Inches(2.6), Inches(0.45),
         PLUM if i < 2 else RGBColor(0xFF,0xFF,0xFF))
    tb(s, x + Inches(0.6), Inches(4.5), Inches(2.6), Inches(0.45),
       sub, 11, bold=True,
       color=CORAL if i < 2 else CORAL, align=PP_ALIGN.CENTER)

tb(s, Inches(0.8), Inches(5.9), Inches(11.7), Inches(0.9),
   "Recovery is the differentiator. A hard week is normal and manageable. Burnout happens when the hard week becomes the normal week, and recovery never catches up.",
   12.5, bold=True, color=BLUSH, align=PP_ALIGN.CENTER, italic=True)

slide_num_label(s)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — WHY TIME OFF DOESN'T WORK
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, CREAM)
top_bar(s, CORAL)

# Right photo
pic(s, IMG + "s06-vacation.png", Inches(7.0), Inches(0.15), Inches(6.33), H - Inches(0.15))
rect(s, Inches(7.0), Inches(0.15), Inches(1.5), H - Inches(0.15), CREAM)  # blend edge

eyebrow(s, "Common Misconception", color=PLUM)
tb(s, Inches(0.8), Inches(1.1), Inches(6.5), Inches(0.85),
   "Why \"Time Off\" Doesn't Always Work", 24, bold=True, color=PLUM)

tb(s, Inches(0.8), Inches(2.1), Inches(6.2), Inches(0.7),
   "The Paradox of Vacation", 18, bold=True, color=CORAL, italic=True)

tb(s, Inches(0.8), Inches(2.85), Inches(6.2), Inches(0.8),
   "You can stop working and still not recover. You can take a week off and return just as exhausted. Why?",
   13, color=CHARCOAL)

reasons = [
    ("Recovery must match depletion type", "Mental exhaustion needs mental rest — not just physical downtime"),
    ("The nervous system stays activated", "Time off doesn't automatically guarantee relaxation or recovery"),
    ("Wrong rest = lingering exhaustion", "Mismatch between rest type and depletion prevents true restoration"),
]
for i, (title, desc) in enumerate(reasons):
    y = Inches(3.8 + i * 1.0)
    rect(s, Inches(0.8), y, Inches(0.06), Inches(0.75), CORAL)
    tb(s, Inches(1.0), y, Inches(5.5), Inches(0.4), title, 13, bold=True, color=PLUM)
    tb(s, Inches(1.0), y + Inches(0.38), Inches(5.5), Inches(0.45), desc, 11.5, color=MUTED)

slide_num_label(s)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — WHAT KIND OF TIRED ARE YOU?
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, PLUM)
side_bar(s)

eyebrow(s, "The Turning Point")
tb(s, Inches(0.6), Inches(1.0), Inches(12.0), Inches(0.9),
   "What Kind of Tired Are You?", 34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

divider_line(s, Inches(5.5), Inches(2.05), Inches(2.33), WHITE)

steps = [
    ("01", "Recognize: Not All Tired Is the Same",
     "Physical exhaustion feels different from mental fog, which differs from emotional depletion or social drain."),
    ("02", "Identify Your Specific Rest Deficit",
     "Pinpoint which dimension is most depleted — body, mind, emotions, relationships, senses, creativity, or spirit."),
    ("03", "Choose Targeted Recovery",
     "Match your rest practice to your specific deficit for effective, efficient restoration."),
]
for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.6 + i * 4.2)
    rect(s, x, Inches(2.4), Inches(3.8), Inches(3.2), PLUM2)
    tb(s, x, Inches(2.55), Inches(3.8), Inches(0.7),
       num, 28, bold=True, color=CORAL, align=PP_ALIGN.CENTER)
    divider_line(s, x + Inches(0.6), Inches(3.25), Inches(2.6), CORAL)
    tb(s, x + Inches(0.15), Inches(3.45), Inches(3.5), Inches(0.6),
       title, 13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, x + Inches(0.15), Inches(4.1), Inches(3.5), Inches(1.2),
       desc, 11.5, color=PALE, align=PP_ALIGN.CENTER)

tb(s, Inches(0.6), Inches(5.9), Inches(12.1), Inches(0.75),
   'Instead of asking “How do I get more rest?” — ask “What kind of rest am I missing?” That question is the turning point.',
   13, bold=True, color=BLUSH, italic=True, align=PP_ALIGN.CENTER)

slide_num_label(s)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — THE 7 TYPES OVERVIEW
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, WHITE)
top_bar(s, PLUM)

eyebrow(s, "The Framework", color=PLUM)
tb(s, Inches(1), Inches(1.1), Inches(11.33), Inches(0.75),
   "The 7 Types of Rest®", 30, bold=True, color=PLUM, align=PP_ALIGN.CENTER)
tb(s, Inches(2), Inches(1.95), Inches(9.33), Inches(0.55),
   "As we explore each type, notice which makes you think: \"That's me.\"",
   14, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

types = [
    ("🛌", "Physical"),
    ("🧠", "Mental"),
    ("💛", "Emotional"),
    ("💬", "Social"),
    ("🔇", "Sensory"),
    ("🎨", "Creative"),
    ("🌿", "Spiritual"),
]
cols = 4
for i, (emoji, name) in enumerate(types):
    col = i % cols
    row = i // cols
    x = Inches(0.9 + col * 3.1)
    y = Inches(2.7 + row * 2.1)
    rect(s, x, y, Inches(2.7), Inches(1.75),
         PLUM if i % 2 == 0 else PLUM2)
    tb(s, x, y + Inches(0.15), Inches(2.7), Inches(0.7),
       emoji, 30, align=PP_ALIGN.CENTER)
    tb(s, x, y + Inches(0.9), Inches(2.7), Inches(0.5),
       f"{i+1}. {name} Rest", 14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

slide_num_label(s)


# ══════════════════════════════════════════════════════════════════════════════
# REST TYPE SLIDES 9-15
# ══════════════════════════════════════════════════════════════════════════════

REST_TYPES = [
    {
        "num": 1, "emoji": "🛌", "name": "Physical",
        "subtitle": "Beyond Sleep",
        "desc": "Physical rest includes sleep, but also recovery that releases physical strain — especially important for desk work, high stress, or bodies carrying chronic tension. If you wake up tired, your body may still be in \"go-mode.\"",
        "col1_title": "Passive Rest", "col1": ["Sleep and naps", "Lying down and resting", "Intentional downtime"],
        "col2_title": "Active Rest", "col2": ["Stretching and gentle movement", "Yoga or tai chi", "Massage or bodywork", "Restorative exercise"],
        "signs": "Persistent fatigue, muscle tension, body aches, low stamina, waking up tired despite adequate sleep hours",
        "photo": None,
        "bg": PLUM,
    },
    {
        "num": 2, "emoji": "🧠", "name": "Mental",
        "subtitle": "Decision Fatigue & Cognitive Load",
        "desc": "The mind that won't turn off. Constant processing, overthinking, and mental loops that continue even when you're trying to relax or sleep.",
        "col1_title": "What It Looks Like",
        "col1": ["Brain fog and forgetfulness", "Difficulty concentrating", "Decision fatigue", "Feeling mentally \"full\""],
        "col2_title": "What You Need",
        "col2": ["Strategic pauses every 2 hours", "Brain dumps to offload looping thoughts", "Reduced context-switching", "Clear shutdown rituals"],
        "signs": "Constant internal chatter, forgetting simple things, irritability, lying awake problem-solving",
        "photo": IMG + "s10-mental.png",
        "bg": PLUM,
    },
    {
        "num": 3, "emoji": "💛", "name": "Emotional",
        "subtitle": "Stop Performing",
        "desc": "When you're constantly managing others' perceptions, you spend enormous emotional energy. Emotional rest is what you desperately need when \"I'm fine\" has become your default response.",
        "col1_title": "Signs of Depletion",
        "col1": ["Chronic people-pleasing", "Irritability and resentment", "Inability to say no", "Emotional numbness"],
        "col2_title": "What Restores You",
        "col2": ["Safe spaces to be honest", "Stop performing \"okay\"", "Set one real boundary", "Express without editing"],
        "signs": "Feeling like you're performing constantly, unexplained resentment, inability to say no without guilt",
        "photo": None,
        "bg": RGBColor(0x3A, 0x2E, 0x5A),
    },
    {
        "num": 4, "emoji": "💬", "name": "Social",
        "subtitle": "People Who Restore vs. Drain",
        "desc": "You can be surrounded by people all day and still feel profoundly isolated or completely drained. Social rest is about the quality of connection, not the quantity.",
        "col1_title": "Warning Signs",
        "col1": ["Feeling lonely despite constant interaction", "Dreading social obligations", "Craving solitude after events", "Depletion after certain relationships"],
        "col2_title": "What Restores You",
        "col2": ["Intentionally choose restorative relationships", "Set limits on draining interactions", "Find authentic community", "Allow yourself to be real"],
        "signs": "Exhausted after interactions you \"should\" enjoy, avoiding previously loved activities",
        "photo": IMG + "s12-social.png",
        "bg": PLUM,
    },
    {
        "num": 5, "emoji": "🔇", "name": "Sensory",
        "subtitle": "Nervous System Recovery",
        "desc": "Screens, notifications, bright lights, background noise, crowded spaces — modern life bombards your nervous system with relentless input it was never designed to handle.",
        "col1_title": "Signs You Need It",
        "col1": ["Frequent headaches", "Feeling \"wired but tired\"", "Light or sound sensitivity", "Agitation in public spaces"],
        "col2_title": "What Helps",
        "col2": ["End screen time 60 min before bed", "Create daily quiet periods", "Dim lighting in evenings", "Eliminate unnecessary notifications"],
        "signs": "Hypersensitivity, overwhelm in crowds, difficulty falling asleep despite exhaustion",
        "photo": IMG + "s13-sensory.png",
        "bg": PLUM,
    },
    {
        "num": 6, "emoji": "🎨", "name": "Creative",
        "subtitle": "Inspiration, Wonder & Problem-Solving",
        "desc": "Creative rest isn't only for artists. It fuels innovation in every role. If everything feels dull, thinking becomes mechanical, or you've lost your sense of wonder — creative rest is the missing piece.",
        "col1_title": "Signs of Deficit",
        "col1": ["Feeling stuck or uninspired", "Cynicism about work and life", "Repetitive, mechanical thinking", "Loss of enthusiasm for projects"],
        "col2_title": "What Restores You",
        "col2": ["Experience beauty — art, nature, music", "Let yourself play without a goal", "Change your environment", "Consume without an agenda"],
        "signs": "Everything feels like a chore; you've lost curiosity about things that used to excite you",
        "photo": IMG + "s14-creative.png",
        "bg": RGBColor(0x3A, 0x2E, 0x5A),
    },
    {
        "num": 7, "emoji": "🌿", "name": "Spiritual",
        "subtitle": "Meaning, Belonging & Purpose",
        "desc": "Spiritual rest is about being anchored — connected to meaning, purpose, and belonging that extends beyond daily tasks. Not necessarily religious, but deeply human.",
        "col1_title": "Signs of Depletion",
        "col1": ["Feeling empty or going through motions", "Questioning \"Why am I doing this?\"", "Loss of meaning in work", "Disconnection from values"],
        "col2_title": "What Restores You",
        "col2": ["Reconnect with core values", "Engage in community or service", "Faith, meditation, or reflection", "Purpose-driven work and impact"],
        "signs": "Even success feels hollow; daily activity feels disconnected from what you believe matters",
        "photo": None,
        "bg": PLUM,
    },
]

for rt in REST_TYPES:
    s = new_slide()
    bg(s, rt["bg"])
    side_bar(s)

    # Slide number badge
    tb(s, W - Inches(2.5), Inches(0.2), Inches(2.2), Inches(0.4),
       f"TYPE {rt['num']} OF 7", 10, bold=True, color=MUTED, align=PP_ALIGN.RIGHT)

    # Emoji + title
    tb(s, Inches(0.5), Inches(0.55), Inches(1.2), Inches(1.0), rt["emoji"], 42)
    tb(s, Inches(1.8), Inches(0.6), Inches(10), Inches(0.7),
       f"{rt['name']} Rest", 28, bold=True, color=WHITE)
    tb(s, Inches(1.8), Inches(1.35), Inches(10), Inches(0.45),
       rt["subtitle"], 14, italic=True, color=CORAL)

    divider_line(s, Inches(0.5), Inches(2.0), Inches(12.33), CORAL)

    if rt["photo"]:
        # Layout: text left, photo right
        pic(s, rt["photo"], Inches(8.5), Inches(2.1), Inches(4.6), Inches(4.9))
        rect(s, Inches(8.5), Inches(2.1), Inches(1.0), Inches(4.9), rt["bg"])  # edge blend
        text_w = Inches(8.0)
    else:
        text_w = Inches(12.4)

    # Description
    tb(s, Inches(0.5), Inches(2.15), text_w, Inches(0.85),
       rt["desc"], 13, color=PALE)

    # Two columns
    for col_i, (c_title, c_items) in enumerate([
        (rt["col1_title"], rt["col1"]),
        (rt["col2_title"], rt["col2"]),
    ]):
        cx = Inches(0.5 + col_i * 4.2) if not rt["photo"] else Inches(0.5 + col_i * 4.0)
        cw = Inches(3.9)
        tb(s, cx, Inches(3.25), cw, Inches(0.35),
           c_title.upper(), 9.5, bold=True,
           color=BLUSH if col_i == 0 else CORAL)
        divider_line(s, cx, Inches(3.65), cw - Inches(0.2),
                     BLUSH if col_i == 0 else CORAL, h=Pt(1.5))
        for j, item in enumerate(c_items):
            tb(s, cx, Inches(3.8 + j * 0.65), cw, Inches(0.6),
               f"• {item}", 12, color=WHITE)

    # Signs box at bottom
    rect(s, Inches(0.5), Inches(6.35), Inches(12.33), Inches(0.75),
         RGBColor(0x22, 0x20, 0x3A))
    tb(s, Inches(0.7), Inches(6.42), Inches(1.8), Inches(0.55),
       "SIGNS:", 10, bold=True, color=CORAL)
    tb(s, Inches(2.4), Inches(6.42), Inches(10.1), Inches(0.55),
       rt["signs"], 11, color=PALE, italic=True)

    slide_num_label(s)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 (was 16) — QUICK SELF-CHECK
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, CREAM)
top_bar(s, CORAL)

pic(s, IMG + "s16-selfcheck.png", Inches(7.5), Inches(0.15), Inches(5.83), H - Inches(0.15))
rect(s, Inches(7.5), Inches(0.15), Inches(1.5), H - Inches(0.15), CREAM)

eyebrow(s, "Pause & Reflect", color=PLUM)
tb(s, Inches(0.8), Inches(1.1), Inches(6.5), Inches(1.0),
   "Quick Self-Check\nTake 30 Seconds", 28, bold=True, color=PLUM)

checks = [
    ("1", "Which 1–2 types stood out most?",
     "Notice which rest types immediately resonated — where you thought \"yes, that's exactly what I'm experiencing.\""),
    ("2", "Where do you feel most depleted?",
     "Body, mind, emotions, senses, relationships, creativity, or spirit? Trust your gut reaction."),
]
for i, (num, q, hint) in enumerate(checks):
    y = Inches(2.55 + i * 1.8)
    rect(s, Inches(0.8), y, Inches(6.0), Inches(1.5), PLUM)
    tb(s, Inches(0.95), y + Inches(0.1), Inches(0.4), Inches(0.55),
       num, 20, bold=True, color=CORAL)
    tb(s, Inches(1.5), y + Inches(0.1), Inches(5.1), Inches(0.5),
       q, 14, bold=True, color=WHITE)
    tb(s, Inches(1.5), y + Inches(0.58), Inches(5.1), Inches(0.75),
       hint, 11.5, color=PALE)

tb(s, Inches(0.8), Inches(6.2), Inches(6.5), Inches(0.7),
   "The types that grabbed you — or the ones you felt resistant to — often reveal your biggest deficits.",
   12, italic=True, color=CORAL)

slide_num_label(s)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — REST MISTAKES HIGH PERFORMERS MAKE
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, PLUM)
side_bar(s)

eyebrow(s, "Common Pitfalls")
tb(s, Inches(0.6), Inches(1.0), Inches(12.1), Inches(0.8),
   "Rest Mistakes High Performers Make", 28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

mistakes = [
    ("1", "Defaulting to Only Physical Rest",
     "Assuming sleep will fix everything, when mental, emotional, or sensory depletion requires entirely different interventions."),
    ("2", "Using Entertainment as Recovery",
     "Scrolling social media, binge-watching, or gaming while remaining overstimulated — mistaking distraction for actual rest."),
    ("3", "Recovering After the Crash",
     "Waiting until complete burnout before addressing rest needs. Prevention requires far less energy than crisis recovery."),
    ("4", "Treating Boundaries as Optional",
     "Believing rest is something you \"earn\" rather than recognizing it as essential infrastructure for sustainable performance."),
]
for i, (num, title, desc) in enumerate(mistakes):
    col = i % 2
    row = i // 2
    x = Inches(0.5 + col * 6.3)
    y = Inches(2.15 + row * 2.1)
    rect(s, x, y, Inches(5.9), Inches(1.85), PLUM2)
    tb(s, x + Inches(0.15), y + Inches(0.1), Inches(0.45), Inches(0.55),
       num, 22, bold=True, color=CORAL)
    tb(s, x + Inches(0.7), y + Inches(0.1), Inches(4.9), Inches(0.45),
       title, 14, bold=True, color=WHITE)
    tb(s, x + Inches(0.7), y + Inches(0.55), Inches(4.9), Inches(1.0),
       desc, 11.5, color=PALE)

tb(s, Inches(0.6), Inches(6.55), Inches(12.1), Inches(0.65),
   "High performers delay rest until they \"earn it.\" But recovery isn't a reward — it's a requirement.",
   13, bold=True, color=CORAL, italic=True, align=PP_ALIGN.CENTER)

slide_num_label(s)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — MICRO-REST / BUILD RESILIENCE
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, WHITE)
top_bar(s, PLUM)

eyebrow(s, "The Solution", color=PLUM)
tb(s, Inches(1), Inches(1.1), Inches(11.33), Inches(0.75),
   "Build Resilience with Micro-Rest", 30, bold=True, color=PLUM, align=PP_ALIGN.CENTER)

tb(s, Inches(1.5), Inches(2.0), Inches(10.33), Inches(0.65),
   "Sustainable renewal doesn't require a life overhaul. It requires small, targeted practices — done consistently — so recovery keeps pace with demand before deficits accumulate.",
   13.5, color=CHARCOAL, align=PP_ALIGN.CENTER)

steps = [
    ("01", "Small Practices", "Manageable actions that fit your real schedule — not an idealized one"),
    ("02", "Repeated Consistently", "Frequency matters more than duration. Daily beats weekly every time"),
    ("03", "Match Rest to Deficit", "Choose the right type for your specific depletion — not a generic fix"),
    ("04", "Protect Like a Meeting", "Schedule recovery with the same priority as your most important work"),
]
for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.6 + i * 3.15)
    rect(s, x, Inches(2.95), Inches(2.85), Inches(3.5), PLUM)
    tb(s, x, Inches(3.1), Inches(2.85), Inches(0.7),
       num, 26, bold=True, color=CORAL, align=PP_ALIGN.CENTER)
    divider_line(s, x + Inches(0.5), Inches(3.85), Inches(1.85), CORAL)
    tb(s, x + Inches(0.15), Inches(4.05), Inches(2.55), Inches(0.5),
       title, 14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, x + Inches(0.15), Inches(4.6), Inches(2.55), Inches(1.5),
       desc, 11.5, color=PALE, align=PP_ALIGN.CENTER)

slide_num_label(s)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — 24-HOUR RESET ACTION PLAN
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, PLUM)
side_bar(s)

eyebrow(s, "Your Action Plan")
tb(s, Inches(0.6), Inches(1.0), Inches(12.0), Inches(0.8),
   "Your 24-Hour Reset", 30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, Inches(1.5), Inches(1.95), Inches(10.33), Inches(0.5),
   "Choose ONE small action to complete in the next 24 hours.", 14, italic=True,
   color=BLUSH, align=PP_ALIGN.CENTER)

actions = [
    ("🛌", "Physical",  "Go to bed 30–60 minutes earlier, or do a 10-minute stretch"),
    ("🧠", "Mental",    "5-minute brain dump, then list tomorrow's top 3 priorities"),
    ("💛", "Emotional", "Send one honest message or set one real boundary"),
    ("💬", "Social",    "Schedule one nourishing connection or cancel a draining obligation"),
    ("🔇", "Sensory",   "30 minutes — no screens, in silence or dim light"),
    ("🎨", "Creative",  "15 minutes of beauty: nature walk, art, or your favorite music"),
    ("🌿", "Spiritual", "10 minutes of prayer, meditation, or reflection on what matters"),
]
cols = 4
for i, (emoji, rest_type, action) in enumerate(actions):
    col = i % cols
    row = i // cols
    x = Inches(0.5 + col * 3.2)
    y = Inches(2.65 + row * 2.05)
    rect(s, x, y, Inches(3.0), Inches(1.8), PLUM2)
    tb(s, x, y + Inches(0.08), Inches(3.0), Inches(0.6),
       emoji, 26, align=PP_ALIGN.CENTER)
    tb(s, x, y + Inches(0.65), Inches(3.0), Inches(0.38),
       rest_type, 12, bold=True, color=CORAL, align=PP_ALIGN.CENTER)
    tb(s, x + Inches(0.1), y + Inches(1.05), Inches(2.8), Inches(0.65),
       action, 10, color=WHITE, align=PP_ALIGN.CENTER)

slide_num_label(s)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — 7-DAY CHALLENGE
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, CREAM)
top_bar(s, CORAL)

eyebrow(s, "7-Day Challenge", color=PLUM)
tb(s, Inches(1), Inches(1.1), Inches(11.33), Inches(0.75),
   "Your Support Plan: Make It Repeatable", 28, bold=True, color=PLUM, align=PP_ALIGN.CENTER)

steps7 = [
    ("1", "Choose 1 Practice",
     "Select one micro-action or ritual that feels achievable and genuinely beneficial for your specific deficit."),
    ("2", "Schedule 3–5 Times This Week",
     "Put it in your calendar as a non-negotiable appointment — treat it like your most important meeting."),
    ("3", "Make It Easy, Not Impressive",
     "Focus on consistent, easy engagement. Sustainable practice yields far better long-term results than heroic one-offs."),
]
for i, (num, title, desc) in enumerate(steps7):
    x = Inches(0.8 + i * 4.1)
    rect(s, x, Inches(2.2), Inches(3.8), Inches(3.5), PLUM)
    tb(s, x, Inches(2.35), Inches(3.8), Inches(0.7),
       num, 32, bold=True, color=CORAL, align=PP_ALIGN.CENTER)
    divider_line(s, x + Inches(0.5), Inches(3.1), Inches(2.8), CORAL)
    tb(s, x + Inches(0.15), Inches(3.3), Inches(3.5), Inches(0.55),
       title, 14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, x + Inches(0.15), Inches(3.95), Inches(3.5), Inches(1.4),
       desc, 12, color=PALE, align=PP_ALIGN.CENTER)

tb(s, Inches(1), Inches(6.0), Inches(11.33), Inches(0.5),
   "If it's not scheduled, it will be skipped. Treat your rest practice like any important appointment.",
   12.5, italic=True, color=CHARCOAL, align=PP_ALIGN.CENTER)

tb(s, Inches(1), Inches(6.6), Inches(11.33), Inches(0.5),
   "Small, consistent actions compound — leading to sustainable renewal.",
   13, bold=True, color=CORAL, align=PP_ALIGN.CENTER)

slide_num_label(s)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 21 — THE REAL CAUSE OF EXHAUSTION (summary)
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, WHITE)
top_bar(s, PLUM)

pic(s, IMG + "s21-summary.png", Inches(7.3), Inches(0.15), Inches(6.03), H - Inches(0.15))
rect(s, Inches(7.3), Inches(0.15), Inches(1.3), H - Inches(0.15), WHITE)

eyebrow(s, "The Bottom Line", color=PLUM)
tb(s, Inches(0.8), Inches(1.1), Inches(6.4), Inches(0.8),
   "The Real Cause\nof Exhaustion", 28, bold=True, color=PLUM)

summaries = [
    ("Exhaustion is a Rest Deficit",
     "You're not lazy, weak, or flawed. You're experiencing a predictable response to chronic depletion without adequate, targeted recovery."),
    ("Resilience Grows with Recovery",
     "True resilience isn't built by enduring more stress — it grows when recovery consistently matches demand."),
    ("Use the 7 Types of Rest®",
     "Identify your specific deficit and choose targeted renewal practices that address your actual depletion — not a generic fix."),
]
for i, (title, desc) in enumerate(summaries):
    y = Inches(2.3 + i * 1.5)
    rect(s, Inches(0.8), y, Inches(0.06), Inches(1.1), CORAL)
    tb(s, Inches(1.1), y, Inches(6.0), Inches(0.45), title, 14, bold=True, color=PLUM)
    tb(s, Inches(1.1), y + Inches(0.45), Inches(6.0), Inches(0.75), desc, 12, color=CHARCOAL)

tb(s, Inches(0.8), Inches(6.95), Inches(6.5), Inches(0.4),
   "When you identify your deficit and practice targeted recovery, resilience becomes something you build — not something you force.",
   10.5, italic=True, color=MUTED)

slide_num_label(s)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 22 — Q&A
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, PLUM)
side_bar(s)

pic(s, IMG + "s22-qa.png", Inches(7.5), 0, Inches(5.83), H)
rect(s, Inches(7.5), 0, Inches(1.3), H, PLUM)

eyebrow(s, "Discussion")
tb(s, Inches(0.6), Inches(1.0), Inches(7.0), Inches(0.85),
   "Q&A + Your Next Step", 30, bold=True, color=WHITE)
tb(s, Inches(0.6), Inches(1.95), Inches(7.0), Inches(0.7),
   "Small actions compound. Targeted rest is how sustainable performance is maintained. Your exhaustion is real — and so is the path forward.",
   12, italic=True, color=PALE)

qs = [
    ("What stood out most?",
     "Which rest type or insight resonated most strongly with your current experience?"),
    ("What is your 24-hour reset?",
     "Which single action will you commit to completing in the next day?"),
    ("What will you repeat for 7 days?",
     "Which practice will you schedule 3–5 times this week to begin building consistency?"),
]
for i, (q, hint) in enumerate(qs):
    y = Inches(3.0 + i * 1.3)
    rect(s, Inches(0.6), y, Inches(6.7), Inches(1.1), PLUM2)
    tb(s, Inches(0.8), y + Inches(0.05), Inches(6.3), Inches(0.45),
       q, 14, bold=True, color=CORAL)
    tb(s, Inches(0.8), y + Inches(0.5), Inches(6.3), Inches(0.5),
       hint, 11.5, color=PALE)

slide_num_label(s)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 23 — THANK YOU / CONTACT
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, PLUM)

side_bar(s)
rect(s, 0, H - Inches(0.15), W, Inches(0.15), CORAL)

# Cindy photo right
pic(s, CINDY_ABOUT, Inches(8.8), 0, Inches(4.53), H)
circ2 = s.shapes.add_shape(9, Inches(9.5), Inches(4.5), Inches(5), Inches(5))
circ2.fill.solid(); circ2.fill.fore_color.rgb = PLUM2; circ2.line.fill.background()

tb(s, Inches(0.7), Inches(1.0), Inches(8), Inches(0.8),
   "Thank You", 42, bold=True, color=WHITE)
tb(s, Inches(0.7), Inches(1.95), Inches(8), Inches(0.5),
   "Cindy Barnard, BSN, RN, OCN", 18, bold=True, color=CORAL)
tb(s, Inches(0.7), Inches(2.5), Inches(8), Inches(0.4),
   "Certified 7 Types of Rest® Facilitator  ·  CIONS Precision Medicine Nursing",
   12, color=PALE, italic=True)

divider_line(s, Inches(0.7), Inches(3.15), Inches(4.5), CORAL)

contact_items = [
    ("✉", "cindy.callingallnurses@gmail.com"),
    ("🌐", "nursecindycares.com"),
    ("📧", "info@nursecindy.com"),
]
for i, (icon, val) in enumerate(contact_items):
    y = Inches(3.45 + i * 0.6)
    tb(s, Inches(0.7), y, Inches(0.5), Inches(0.5), icon, 15)
    tb(s, Inches(1.3), y + Inches(0.06), Inches(7.0), Inches(0.45), val, 14, color=WHITE)

tb(s, Inches(0.7), Inches(5.5), Inches(8.0), Inches(0.7),
   "Discover Your Rest Deficit", 16, bold=True, color=BLUSH)
tb(s, Inches(0.7), Inches(6.1), Inches(8.0), Inches(0.7),
   "Visit RestQuiz.com for the free rest deficit assessment\nby Dr. Saundra Dalton-Smith.", 12, color=PALE)

tb(s, Inches(0.7), Inches(7.0), Inches(8.0), Inches(0.35),
   "7 Types of Rest® is a registered trademark of Dr. Saundra Dalton-Smith / Restorasis.",
   9, color=MUTED, italic=True)

slide_num_label(s, 23)


# ── Save ───────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved → {OUT}")
print(f"Slides: {len(prs.slides)}")
