"""Build Nurse Cindy presentation deck using python-pptx."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Brand colors ──────────────────────────────────────────────────────────────
PLUM   = RGBColor(0x2C, 0x2A, 0x4A)
CORAL  = RGBColor(0xE8, 0x89, 0x6A)
CREAM  = RGBColor(0xFA, 0xF7, 0xF4)
BLUSH  = RGBColor(0xF2, 0xD4, 0xCC)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
CHARCOAL = RGBColor(0x3D, 0x3D, 0x3D)
MUTED  = RGBColor(0x6B, 0x6B, 0x6B)

W  = Inches(13.33)   # widescreen 16:9 width
H  = Inches(7.5)

OUT = os.path.join(os.path.dirname(__file__), "assets", "NurseCindy-Deck.pptx")

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # completely blank

# ── Helpers ───────────────────────────────────────────────────────────────────

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, color, alpha=None):
    shp = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp

def txbox(slide, l, t, w, h, text, size, bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, italic=False, font="Calibri"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return tb

def eyebrow(slide, text, top=Inches(1.2)):
    txbox(slide, Inches(1), top, Inches(11.33), Inches(0.4),
          text.upper(), 11, bold=True, color=CORAL, align=PP_ALIGN.CENTER)

def divider(slide, l, t, w=Inches(1), color=CORAL):
    rect(slide, l, t, w, Pt(3), color)

# ── Slide 1: Title ─────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, PLUM)

# Coral accent bar left
rect(s, 0, 0, Inches(0.25), H, CORAL)

# Coral accent circle (decorative)
circ = s.shapes.add_shape(9, Inches(9.5), Inches(-1), Inches(5), Inches(5))
circ.fill.solid()
circ.fill.fore_color.rgb = RGBColor(0x3D, 0x3A, 0x6E)
circ.line.fill.background()

# Eyebrow
txbox(s, Inches(1.2), Inches(1.8), Inches(8), Inches(0.5),
      "CERTIFIED 7 TYPES OF REST® FACILITATOR", 11, bold=True, color=CORAL)

# Title
txbox(s, Inches(1.2), Inches(2.5), Inches(9), Inches(1.5),
      "You Can't Pour From\nAn Empty Cup.", 44, bold=True, color=WHITE, font="Calibri")

# Subtitle
txbox(s, Inches(1.2), Inches(4.3), Inches(8.5), Inches(0.8),
      "Rest Science for Real People — Keynote & Workshop Series", 18,
      color=BLUSH, italic=True)

# Speaker name
txbox(s, Inches(1.2), Inches(5.4), Inches(6), Inches(0.5),
      "Nurse Cindy  |  nursecindycares.com", 14, color=RGBColor(0xA0, 0x9E, 0xBE))

# Bottom coral line
rect(s, Inches(1.2), Inches(6.8), Inches(5), Pt(2), CORAL)


# ── Slide 2: The Problem ───────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, CREAM)

rect(s, 0, 0, W, Inches(0.18), CORAL)

eyebrow(s, "The Problem", top=Inches(0.6))
txbox(s, Inches(1), Inches(1.3), Inches(11.33), Inches(0.9),
      "We're Exhausted. And Sleep Isn't Fixing It.", 34, bold=True, color=PLUM,
      align=PP_ALIGN.CENTER)

divider(s, Inches(6), Inches(2.5), color=CORAL)

stats = [
    ("77%", "of workers\nexperience burnout"),
    ("67%", "say they sleep enough\nbut still feel tired"),
    ("$322B", "lost annually to\nburnout-related turnover"),
]
for i, (num, label) in enumerate(stats):
    x = Inches(1.2 + i * 3.8)
    rect(s, x, Inches(2.9), Inches(3.3), Inches(2.8), PLUM)
    txbox(s, x, Inches(3.0), Inches(3.3), Inches(1.1),
          num, 42, bold=True, color=CORAL, align=PP_ALIGN.CENTER)
    txbox(s, x, Inches(4.1), Inches(3.3), Inches(0.8),
          label, 14, color=WHITE, align=PP_ALIGN.CENTER)

txbox(s, Inches(1), Inches(6.4), Inches(11.33), Inches(0.6),
      "The cause isn't lack of sleep. It's the wrong kind of rest.", 16,
      bold=True, color=CORAL, align=PP_ALIGN.CENTER, italic=True)


# ── Slide 3: About Nurse Cindy ─────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, PLUM)

# Right blush panel
rect(s, Inches(8.5), 0, Inches(4.83), H, RGBColor(0x3D, 0x3A, 0x6E))

eyebrow(s, "Your Facilitator", top=Inches(0.7))

txbox(s, Inches(0.8), Inches(1.4), Inches(7.2), Inches(1.0),
      "Nurse Cindy", 40, bold=True, color=WHITE)

txbox(s, Inches(0.8), Inches(2.6), Inches(7.0), Inches(0.5),
      "Certified 7 Types of Rest® Facilitator", 16, color=CORAL, italic=True)

divider(s, Inches(0.8), Inches(3.3), Inches(1.5))

creds = [
    ("🏥", "Oncology Nurse", "Frontline clinical experience across oncology and acute care"),
    ("🌐", "Corporate Biotech", "Global industry background in life sciences and healthcare strategy"),
    ("✅", "Certified Facilitator", "Trained under Dr. Saundra Dalton-Smith's research-backed 7 Types of Rest® program"),
]
for i, (icon, title, desc) in enumerate(creds):
    y = Inches(3.7 + i * 1.0)
    txbox(s, Inches(0.8), y, Inches(0.5), Inches(0.5), icon, 18, color=CORAL)
    txbox(s, Inches(1.5), y, Inches(6.0), Inches(0.35),
          title, 14, bold=True, color=WHITE)
    txbox(s, Inches(1.5), y + Inches(0.3), Inches(6.5), Inches(0.5),
          desc, 12, color=RGBColor(0xB0, 0xAE, 0xC8))

txbox(s, Inches(8.7), Inches(2.5), Inches(4.3), Inches(3.5),
      '"I bring a clinical, evidence-based lens to one of the most misunderstood forces in our lives: rest."',
      15, italic=True, color=BLUSH, align=PP_ALIGN.CENTER)

txbox(s, Inches(8.7), Inches(6.0), Inches(4.3), Inches(0.5),
      "— Nurse Cindy", 13, color=CORAL, align=PP_ALIGN.CENTER)


# ── Slide 4: What IS Rest? ─────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
rect(s, 0, 0, W, Inches(0.18), PLUM)

eyebrow(s, "Rest Redefined", top=Inches(0.55))
txbox(s, Inches(1), Inches(1.2), Inches(11.33), Inches(0.8),
      "Rest Is Not Just Sleep", 36, bold=True, color=PLUM, align=PP_ALIGN.CENTER)

txbox(s, Inches(2), Inches(2.2), Inches(9.33), Inches(0.7),
      "Most people treat rest as a single thing — sleep. But research by Dr. Saundra Dalton-Smith\nidentifies SEVEN distinct types of rest. Being deficient in even one leads to exhaustion\nthat no amount of sleep can cure.",
      14, color=CHARCOAL, align=PP_ALIGN.CENTER)

# Two columns
for col, (title, items) in enumerate([
    ("Signs You Need More Rest", [
        "You sleep 7–8 hrs and still feel drained",
        "You're irritable with people you love",
        "Creativity feels blocked or flat",
        "You feel numb or disconnected",
        "Your to-do list never feels done",
    ]),
    ("What This Framework Delivers", [
        "Identify YOUR specific rest deficits",
        "Match rest type to root cause of exhaustion",
        "Build targeted, practical rest strategies",
        "Sustainable — no life overhaul required",
        "Backed by research and clinical experience",
    ]),
]):
    x = Inches(0.8 + col * 6.3)
    rect(s, x, Inches(3.4), Inches(5.8), Inches(0.45), PLUM if col == 0 else CORAL)
    txbox(s, x, Inches(3.4), Inches(5.8), Inches(0.45),
          title, 13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        txbox(s, x + Inches(0.15), Inches(3.95 + j * 0.55), Inches(5.5), Inches(0.5),
              f"• {item}", 13, color=CHARCOAL)


# ── Slides 5–11: One per Rest Type ────────────────────────────────────────────
rest_types = [
    ("🛌", "Physical Rest",
     "The Foundation",
     "Rest for the body — passive (sleep, naps) and active (yoga, stretching, breathwork). Most people stop here. Physical rest is necessary, but not sufficient.",
     ["Persistent physical fatigue", "Muscle tension and soreness", "Poor sleep quality despite hours in bed"],
     ["Prioritize sleep hygiene", "Add gentle movement between intense activity", "Schedule true recovery days — not just less-busy days"]),

    ("🧠", "Mental Rest",
     "Quiet the Mind",
     "Relief for the overworked, decision-fatigued brain. Essential for high performers, managers, caregivers, and anyone who is always the person thinking.",
     ["Difficulty concentrating", "Mental chatter that won't stop", "Forgetting simple things, feeling foggy"],
     ["Set a 2-minute break every 2 hours", "Keep a notepad to offload looping thoughts", "End work with a shutdown ritual to close open loops"]),

    ("💬", "Social Rest",
     "Right Relationships",
     "The ability to differentiate between relationships that restore and those that drain — and to get more of the former. This isn't about being antisocial; it's about being intentional.",
     ["Dread seeing certain people (even ones you love)", "Feel lonely even in groups", "Exhausted after interactions you \"should\" enjoy"],
     ["Identify your 3–5 restorative relationships", "Set limits on energy-draining interactions", "Be real instead of performing — authentic connection restores"]),

    ("🎨", "Creative Rest",
     "Inspire the Imagination",
     "The space that allows wonder, inspiration, and beauty to replenish the parts of us that problem-solve, create, and innovate. Without it, thinking becomes mechanical.",
     ["Everything feels like a chore or obligation", "You've lost curiosity about things that used to excite you", "Problem-solving feels forced and flat"],
     ["Consume beauty — art, nature, music — without an agenda", "Let yourself play without a goal", "Take a different route. Rearrange your environment. Notice more."]),

    ("💛", "Emotional Rest",
     "Permission to Feel",
     "The freedom to express how you actually feel rather than the feelings others need you to have. Critical for caregivers, leaders, parents — anyone who has to manage their emotions for others.",
     ["Performing \"okay\" when you're not", "Feeling resentful without knowing why", "Emotional numbing or unexpected breakdowns"],
     ["Find one person you can be completely honest with", "Stop saying \"I'm fine\" when you're not", "Journal without editing — no performance, just truth"]),

    ("🔇", "Sensory Rest",
     "Turn Down the Volume",
     "Relief from the constant input of screens, noise, notifications, and artificial environments. Our nervous systems were not designed for modern levels of stimulation — and we're paying for it.",
     ["Hypersensitivity to sound, light, or touch", "End-of-day headaches or eye strain", "Feeling overwhelmed in public spaces"],
     ["Deliberately end screen time 60 min before bed", "Create a daily quiet period — even 10 minutes", "Declutter your visual environment"]),

    ("🌿", "Spiritual Rest",
     "Connect to Purpose",
     "The sense that your life and work are part of something larger than yourself. Not necessarily religious — but deeply human. Without it, even success feels hollow.",
     ["Feeling like your work doesn't matter", "Going through the motions without meaning", "Disconnected from what you originally cared about"],
     ["Reconnect with your core values — write them down", "Engage in community or service beyond yourself", "Ask regularly: 'Does this align with what I believe matters?'"]),
]

for emoji, name, subtitle, desc, signs, strategies in rest_types:
    s = prs.slides.add_slide(BLANK)
    bg(s, PLUM)

    # Left coral accent
    rect(s, 0, 0, Inches(0.2), H, CORAL)

    # Type number indicator top right
    idx = rest_types.index((emoji, name, subtitle, desc, signs, strategies)) + 1
    txbox(s, Inches(11), Inches(0.3), Inches(2), Inches(0.5),
          f"Type {idx} of 7", 11, color=RGBColor(0x6B, 0x69, 0x8E), align=PP_ALIGN.RIGHT)

    # Emoji
    txbox(s, Inches(0.5), Inches(0.6), Inches(1.2), Inches(1.0), emoji, 48)

    # Name
    txbox(s, Inches(1.8), Inches(0.6), Inches(9), Inches(0.8),
          f"{name} Rest", 32, bold=True, color=WHITE)
    txbox(s, Inches(1.8), Inches(1.4), Inches(9), Inches(0.45),
          subtitle, 14, color=CORAL, italic=True)

    divider(s, Inches(0.5), Inches(2.1), Inches(12.3), CORAL)

    # Description
    txbox(s, Inches(0.5), Inches(2.3), Inches(12.33), Inches(0.9),
          desc, 14, color=RGBColor(0xC8, 0xC6, 0xD8))

    # Two columns: Signs | Strategies
    for col, (col_title, col_items, col_color) in enumerate([
        ("Signs of Deficit", signs, BLUSH),
        ("Recovery Strategies", strategies, CORAL),
    ]):
        x = Inches(0.5 + col * 6.5)
        txbox(s, x, Inches(3.5), Inches(6.0), Inches(0.4),
              col_title.upper(), 10, bold=True, color=col_color)
        divider(s, x, Inches(3.92), Inches(5.5), col_color)
        for j, item in enumerate(col_items):
            txbox(s, x, Inches(4.05 + j * 0.8), Inches(6.1), Inches(0.7),
                  f"• {item}", 13, color=WHITE)


# ── Slide 12: Speaking Topics ──────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, CREAM)
rect(s, 0, 0, W, Inches(0.18), CORAL)

eyebrow(s, "Speaking Topics", top=Inches(0.55))
txbox(s, Inches(1), Inches(1.2), Inches(11.33), Inches(0.8),
      "Five Signature Talks", 34, bold=True, color=PLUM, align=PP_ALIGN.CENTER)

topics = [
    "Finding Calm in a Constantly Stimulating World",
    "Rest That Fits Real Life",
    "Beyond Burnout: Understanding the Real Causes of Exhaustion",
    "Emotional Skills That Keep Relationships Strong",
    "The Hidden Cost of Always Being the Strong One",
]

for i, topic in enumerate(topics):
    x = Inches(0.9) if i < 3 else Inches(4.8)
    y = Inches(2.2 + (i % 3) * 1.4) if i < 3 else Inches(2.2 + (i - 3) * 1.4)
    if i == 3:
        x = Inches(0.9)
        y = Inches(6.3) if False else Inches(2.2 + (i % 3 if i < 3 else i - 2) * 1.4)

    # simpler: just stack them
    y = Inches(2.1 + i * 0.98)
    shp = rect(s, Inches(0.9), y, Inches(11.5), Inches(0.78),
               PLUM if i % 2 == 0 else RGBColor(0x3D, 0x3A, 0x6E))
    txbox(s, Inches(1.4), y + Inches(0.08), Inches(0.5), Inches(0.6),
          str(i + 1), 20, bold=True, color=CORAL)
    txbox(s, Inches(2.0), y + Inches(0.15), Inches(9.8), Inches(0.5),
          topic, 15, bold=True, color=WHITE)


# ── Slide 13: Engagement Formats ──────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, PLUM)

eyebrow(s, "How We Work Together", top=Inches(0.5))
txbox(s, Inches(1), Inches(1.1), Inches(11.33), Inches(0.8),
      "Flexible Engagement Formats", 34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

formats = [
    ("🎤", "Keynote", "45–90 min", "High-impact, inspiring, and immediately actionable"),
    ("🧩", "Half-Day\nWorkshop", "3–4 hours", "Guided assessment, small groups, action planning"),
    ("📅", "Full-Day\nTraining", "Full day", "Deep-dive for HR, wellness teams, or leadership cohorts"),
    ("💻", "Virtual\nSession", "Any length", "All formats available live online for distributed teams"),
]

for i, (icon, fmt, duration, desc) in enumerate(formats):
    x = Inches(0.6 + i * 3.1)
    rect(s, x, Inches(2.2), Inches(2.8), Inches(4.3), RGBColor(0x3D, 0x3A, 0x6E))
    txbox(s, x, Inches(2.4), Inches(2.8), Inches(0.8), icon, 36, align=PP_ALIGN.CENTER)
    txbox(s, x, Inches(3.2), Inches(2.8), Inches(0.7),
          fmt, 16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x + Inches(0.6), Inches(3.95), Inches(1.6), Pt(2), CORAL)
    txbox(s, x, Inches(4.1), Inches(2.8), Inches(0.4),
          duration, 11, color=CORAL, align=PP_ALIGN.CENTER, bold=True)
    txbox(s, x + Inches(0.1), Inches(4.55), Inches(2.6), Inches(0.9),
          desc, 11, color=RGBColor(0xB0, 0xAE, 0xC8), align=PP_ALIGN.CENTER)

txbox(s, Inches(1), Inches(6.7), Inches(11.33), Inches(0.5),
      "All topics can be customized for your specific audience, industry, and outcomes.",
      12, italic=True, color=CORAL, align=PP_ALIGN.CENTER)


# ── Slide 14: Testimonials ─────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, CREAM)
rect(s, 0, 0, W, Inches(0.18), PLUM)

eyebrow(s, "What People Say", top=Inches(0.5))
txbox(s, Inches(1), Inches(1.1), Inches(11.33), Inches(0.8),
      "Voices From the Room", 34, bold=True, color=PLUM, align=PP_ALIGN.CENTER)

testimonials = [
    ('"Her insights on emotional rest were eye-opening.\nI feel more balanced and at peace."',
     "Jimmelynn Rice", "CEO, Girls Nite In International"),
    ('"Cindy\'s passion for rest is contagious. I left feeling\nempowered and ready to implement."',
     "Jessica Baker", "Advisor-Scientist, Eli Lilly & Co."),
    ('"Nurse Cindy\'s talk transformed my perspective on rest.\nThis showed me how wrong I was — in the best way."',
     "Cindy Modafferi, RN", "Registered Nurse"),
]

for i, (quote, name, title) in enumerate(testimonials):
    x = Inches(0.5 + i * 4.3)
    rect(s, x, Inches(2.2), Inches(3.9), Inches(4.5), PLUM)
    txbox(s, x + Inches(0.15), Inches(2.3), Inches(0.6), Inches(0.7),
          "“", 52, color=CORAL)
    txbox(s, x + Inches(0.25), Inches(3.0), Inches(3.5), Inches(2.0),
          quote, 12, italic=True, color=WHITE)
    divider(s, x + Inches(0.25), Inches(5.1), Inches(2.5), CORAL)
    txbox(s, x + Inches(0.25), Inches(5.25), Inches(3.5), Inches(0.4),
          name, 13, bold=True, color=BLUSH)
    txbox(s, x + Inches(0.25), Inches(5.6), Inches(3.5), Inches(0.4),
          title, 11, color=RGBColor(0xA0, 0x9E, 0xBE))


# ── Slide 15: How to Book / CTA ────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s, PLUM)

# Coral accent
rect(s, 0, 0, Inches(0.25), H, CORAL)
rect(s, 0, H - Inches(0.25), W, Inches(0.25), CORAL)

txbox(s, Inches(1.2), Inches(1.2), Inches(11), Inches(0.5),
      "READY TO BRING REST SCIENCE TO YOUR TEAM?", 13, bold=True, color=CORAL,
      align=PP_ALIGN.CENTER)

txbox(s, Inches(1.2), Inches(1.9), Inches(11), Inches(1.3),
      "Book Nurse Cindy\nfor Your Next Event", 38, bold=True, color=WHITE,
      align=PP_ALIGN.CENTER)

txbox(s, Inches(2.5), Inches(3.4), Inches(8.33), Inches(0.7),
      "Nurse Cindy works with organizations of all sizes — from intimate leadership retreats\nto large conferences — to deliver experiences that change how people relate to rest.",
      14, color=RGBColor(0xB0, 0xAE, 0xC8), align=PP_ALIGN.CENTER)

contact_items = [
    ("✉", "Email", "cindy.callingallnurses@gmail.com"),
    ("🌐", "Website", "nursecindycares.com"),
    ("📧", "General", "info@nursecindy.com"),
]
for i, (icon, label, value) in enumerate(contact_items):
    x = Inches(1.5 + i * 3.8)
    rect(s, x, Inches(4.4), Inches(3.3), Inches(1.4), RGBColor(0x3D, 0x3A, 0x6E))
    txbox(s, x, Inches(4.5), Inches(3.3), Inches(0.45),
          f"{icon}  {label}", 11, bold=True, color=CORAL, align=PP_ALIGN.CENTER)
    txbox(s, x, Inches(4.95), Inches(3.3), Inches(0.5),
          value, 12, color=WHITE, align=PP_ALIGN.CENTER)

txbox(s, Inches(1), Inches(6.3), Inches(11.33), Inches(0.7),
      "Empowered by Rest.  |  nursecindycares.com",
      13, color=RGBColor(0x6B, 0x69, 0x8E), align=PP_ALIGN.CENTER)

# ── Save ───────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"Saved: {OUT}")
